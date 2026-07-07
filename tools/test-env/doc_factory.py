"""テスト文書ファクトリ。

org_model.FOLDERS の doc_category 毎に、多様な内容の日本語ビジネス文書を
Office 形式 (docx / xlsx / pptx / pdf) で決定論的に生成する。

設計方針:
  - 乱数は random.Random(f"{category}:{SEED}") で決定論化 (再実行で同一内容)
  - ベクトル検索でエリア毎の差が際立つよう、カテゴリ固有の語彙を濃く入れる
    (例: 新製品X のコードネーム「Aurora」は projx_* と取締役会議事録にしか出ない)
  - pptx は NemakiWare の rag.supported.mimetypes 既定に含まれず RAG 索引外。
    シナリオで参照したい内容は必ず docx/xlsx/pdf 側にも存在させる
  - PDF は日本語フォントが見つからない環境では docx にフォールバック
"""

from __future__ import annotations

import io
import random
from dataclasses import dataclass
from typing import Callable, Iterator

SEED = 20260707

# ---------------------------------------------------------------------------
# 素材プール (架空の会社・顧客・製品)
# ---------------------------------------------------------------------------
COMPANY = "ヒナタ産業株式会社"

PRODUCTS = {
    "DocuHive": "クラウド文書管理サービス",
    "SignFlow": "電子契約ワークフローサービス",
    "InsightPad": "データ分析ダッシュボード",
}
# 新製品X (機密): projx_* 文書と取締役会議事録にのみ登場させる
PROJX_CODENAME = "Aurora"
PROJX_DESC = "AI文書検索アシスタント"

CUSTOMERS = [
    "株式会社トチギ電機", "北都銀行", "ミナモト製薬", "株式会社ヤマセ運輸",
    "つばき生命保険", "株式会社コダマ建設", "エミシア食品株式会社", "南海ガス",
    "株式会社フジサキ商店", "ハルカゼ大学", "株式会社ミドリ精機", "青葉市役所",
    "株式会社ソラノ通信", "タチバナ不動産", "株式会社クロベ鉄鋼", "ワカバ学園",
    "株式会社ナルミ電子", "ヒカリ物流株式会社", "株式会社アカツキ印刷", "セイラン化学株式会社",
]

REGIONS = ["東京", "大阪", "名古屋", "福岡", "札幌", "仙台", "広島", "横浜"]

SALES_STAFF = ["森 大輔", "浅田 未来", "上田 修", "星野 佳奈", "工藤 誠"]
DEV_STAFF = ["福田 亮", "小川 蓮", "西田 樹", "宮田 聡", "平田 淳"]

MONTHS = [
    "2025年10月", "2025年11月", "2025年12月",
    "2026年1月", "2026年2月", "2026年3月", "2026年4月", "2026年5月", "2026年6月",
]

MIME = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf": "application/pdf",
}


@dataclass(frozen=True)
class GeneratedDoc:
    category: str
    filename: str
    mimetype: str
    content: bytes


# ---------------------------------------------------------------------------
# 低レベルビルダー
# ---------------------------------------------------------------------------

def make_docx(title: str, meta_lines: list[str], sections: list[tuple[str, list[str]]]) -> bytes:
    from docx import Document

    doc = Document()
    doc.add_heading(title, level=0)
    for line in meta_lines:
        doc.add_paragraph(line)
    for heading, paragraphs in sections:
        doc.add_heading(heading, level=1)
        for p in paragraphs:
            doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_xlsx(sheet_title: str, title: str, meta_lines: list[str],
              headers: list[str], rows: list[list], notes: list[str]) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = sheet_title[:28]
    ws["A1"] = title
    ws["A1"].font = Font(bold=True, size=14)
    r = 3
    for line in meta_lines:
        ws.cell(row=r, column=1, value=line)
        r += 1
    r += 1
    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF")
    for c, h in enumerate(headers, 1):
        cell = ws.cell(row=r, column=c, value=h)
        cell.fill = header_fill
        cell.font = header_font
    for row in rows:
        r += 1
        for c, v in enumerate(row, 1):
            cell = ws.cell(row=r, column=c, value=v)
            if isinstance(v, int):
                cell.number_format = "#,##0"
    r += 2
    for note in notes:
        ws.cell(row=r, column=1, value=note)
        r += 1
    for col_idx in range(1, len(headers) + 1):
        ws.column_dimensions[chr(ord("A") + col_idx - 1)].width = 22
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def make_pptx(title: str, subtitle: str, slides: list[tuple[str, list[str]]]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    for slide_title, bullets in slides:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = slide_title
        tf = s.placeholders[1].text_frame
        tf.text = bullets[0] if bullets else ""
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


_JP_FONT_STATE: dict = {}


def _register_jp_font() -> bool:
    """reportlab に日本語フォントを登録する。成功可否をキャッシュ。"""
    if "ok" in _JP_FONT_STATE:
        return _JP_FONT_STATE["ok"]
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    candidates = [
        "/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
        "/System/Library/Fonts/ヒラギノ明朝 ProN.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
    ]
    for path in candidates:
        try:
            pdfmetrics.registerFont(TTFont("JPFont", path, subfontIndex=0))
            _JP_FONT_STATE["ok"] = True
            return True
        except Exception:
            continue
    _JP_FONT_STATE["ok"] = False
    return False


def make_pdf_or_docx(title: str, meta_lines: list[str],
                     sections: list[tuple[str, list[str]]]) -> tuple[bytes, str]:
    """日本語フォントが使えれば PDF、無ければ docx を返す。(bytes, ext)。"""
    if not _register_jp_font():
        return make_docx(title, meta_lines, sections), "docx"

    from reportlab.lib import enums
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    style_body = ParagraphStyle(name="JP", fontName="JPFont", fontSize=10, leading=15)
    style_title = ParagraphStyle(name="JPTitle", fontName="JPFont", fontSize=17,
                                 leading=22, alignment=enums.TA_CENTER, spaceAfter=10)
    style_h1 = ParagraphStyle(name="JPH1", fontName="JPFont", fontSize=13,
                              leading=18, spaceBefore=10, spaceAfter=6)

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, topMargin=22 * mm, bottomMargin=22 * mm,
                            leftMargin=20 * mm, rightMargin=20 * mm)
    story = [Paragraph(title, style_title)]
    for line in meta_lines:
        story.append(Paragraph(line, style_body))
    story.append(Spacer(1, 8 * mm))
    for heading, paragraphs in sections:
        story.append(Paragraph(heading, style_h1))
        for p in paragraphs:
            story.append(Paragraph(p, style_body))
    doc.build(story)
    return buf.getvalue(), "pdf"


# ---------------------------------------------------------------------------
# カテゴリ別ジェネレータ
# 各関数は (rnd, i) -> GeneratedDoc 相当の (filename, mimetype, bytes)
# ---------------------------------------------------------------------------

def gen_rules(rnd: random.Random, i: int):
    topics = [
        ("就業規則", "勤務時間・休日・服務規律", [
            "所定勤務時間は1日8時間、週40時間とする。始業9:00、終業18:00、休憩1時間。",
            "フレックスタイム制のコアタイムは10:00〜15:00とする。",
            "年次有給休暇は勤続年数に応じて年10日から20日を付与する。",
        ]),
        ("賞与規程", "賞与の支給基準と支給時期", [
            "賞与は年2回、6月および12月に支給する。支給額は基本給に評価係数を乗じて算定する。",
            "評価係数は人事評価制度に基づき0.8〜1.5の範囲で決定する。",
            "支給日に在籍しない者には支給しない。",
        ]),
        ("旅費規程", "出張旅費の支給基準", [
            "国内出張の日当は3,000円、海外出張の日当は5,000円とする。",
            "宿泊費は国内上限12,000円、海外上限20,000円を実費支給する。",
            "出張終了後5営業日以内に出張報告書を提出しなければならない。",
        ]),
        ("情報セキュリティ規程", "情報資産の保護", [
            "文書管理システムのアクセス権限は最小権限の原則に基づき部署単位で設定する。",
            "機密文書は指定フォルダ以外への複製を禁止する。",
            "パスワードは12文字以上とし、多要素認証の利用を推奨する。",
        ]),
        ("経費精算規程", "経費の精算手続", [
            "経費精算は発生月の翌月10日までに申請する。",
            "1件5万円以上の支出は事前に所属長の承認を要する。",
            "領収書は電子帳簿保存法の要件に従い電子保存する。",
        ]),
        ("在宅勤務規程", "テレワークの運用", [
            "在宅勤務は週3日を上限とし、前週金曜までに所属長に申請する。",
            "在宅勤務時の通信費補助として月3,000円を支給する。",
            "VPN接続と画面ロックの徹底を義務付ける。",
        ]),
        ("育児介護休業規程", "育児・介護休業の取得", [
            "育児休業は子が2歳に達するまで取得できる。",
            "介護休業は対象家族1人につき通算93日まで分割取得できる。",
            "復職後6か月間は短時間勤務を選択できる。",
        ]),
    ]
    name, purpose, articles = topics[i % len(topics)]
    rev = i // len(topics) + 1
    title = f"{name}" + (f"（第{rev}版）" if rev > 1 else "")
    body = [f"第{n}条 {a}" for n, a in enumerate(articles, 1)]
    body.append(f"附則: この規程は{rnd.choice(MONTHS)}1日から施行する。改訂責任部署は人事課とする。")
    content = make_docx(title, [COMPANY, f"文書番号: RULE-{2026}-{i+1:03d}"],
                        [("目的", [f"本規程は{COMPANY}における{purpose}に関する事項を定める。"]),
                         ("本文", body)])
    return f"{title}.docx", MIME["docx"], content


def gen_announcements(rnd: random.Random, i: int):
    items = [
        ("健康診断実施のお知らせ", "定期健康診断を実施します。全従業員は期間内に受診してください。受診票は人事課から配布します。"),
        ("年末年始休業のお知らせ", "12月29日から1月3日まで全社休業とします。緊急連絡は各本部の当番体制に従ってください。"),
        ("オフィス移転について", "本社オフィスを移転します。新オフィスはフリーアドレス制となり、文書は原則電子化して文書管理システムに保存してください。"),
        ("社内表彰式のご案内", "上期の社内表彰式を開催します。営業成績優秀者および業務改善提案の表彰を行います。"),
        ("防災訓練実施のお知らせ", "全社防災訓練を実施します。避難経路の確認と安否確認システムの応答訓練を行います。"),
        ("新入社員歓迎会のご案内", "新入社員の歓迎会を開催します。各部署の紹介と懇親を予定しています。"),
        ("システムメンテナンスのお知らせ", "文書管理システムのメンテナンスを実施します。当日夜間はアクセスできません。"),
        ("インフルエンザ予防接種の補助について", "予防接種費用を1人3,000円まで補助します。領収書を添えて経理課に申請してください。"),
        ("社内公募制度のご案内", "新規事業部門への社内公募を開始します。応募は人事課まで。"),
        ("福利厚生ポイント制度の改定", "カフェテリアプランのポイント付与を年間60,000ポイントに拡充します。"),
        ("駐車場利用ルールの変更", "来客用駐車場の予約制を導入します。総務システムから予約してください。"),
        ("クールビズ実施のお知らせ", "5月から9月までクールビズを実施します。"),
    ]
    title, body = items[i % len(items)]
    month = rnd.choice(MONTHS)
    content = make_docx(f"社内通達: {title}", [COMPANY, f"発信: 人事課 / {month}", f"通達番号: ANN-{i+1:03d}"],
                        [("本文", [body, "不明点は人事課（内線200）までお問い合わせください。"])])
    return f"社内通達_{i+1:02d}_{title}.docx", MIME["docx"], content


def gen_board_minutes(rnd: random.Random, i: int):
    month = MONTHS[i % len(MONTHS)]
    agenda_pool = [
        f"{month}度の連結売上は前年同月比{rnd.randint(95, 118)}%で推移。{rnd.choice(list(PRODUCTS))}の更新率が改善。",
        f"{rnd.choice(REGIONS)}支社の増床について、投資額{rnd.randint(3, 12) * 10}百万円で承認。",
        f"中期経営計画の重点施策「クラウドサービス比率50%」の進捗を確認。現状{rnd.randint(30, 45)}%。",
        f"人材採用計画: エンジニア{rnd.randint(5, 15)}名、営業{rnd.randint(3, 8)}名の増員を決議。",
        f"為替影響によるコスト増{rnd.randint(2, 9)}百万円を織り込み、通期予想を据え置き。",
    ]
    secret = (
        f"新製品X（コードネーム{PROJX_CODENAME}）の開発進捗を確認。{PROJX_DESC}として"
        f"{rnd.choice(['2026年10月', '2026年12月', '2027年1月'])}のリリースを目指す。"
        "本件は新製品Xプロジェクトメンバー以外への情報共有を禁止する。"
    )
    agendas = rnd.sample(agenda_pool, 3) + ([secret] if i % 2 == 0 else [])
    content = make_docx(
        f"取締役会議事録 {month}",
        [COMPANY, f"開催: {month}下旬 / 本社役員会議室", "出席: 大塚社長、工藤営業本部長、平田技術本部長、永井管理本部長"],
        [("決議・報告事項", agendas),
         ("次回", [f"次回取締役会は{rnd.choice(MONTHS)}に開催する。"])])
    return f"取締役会議事録_{month}.docx", MIME["docx"], content


def gen_mid_term_plan(rnd: random.Random, i: int):
    year = 2026 + i % 3
    themes = [
        ("クラウドシフト", "オンプレミス保守収益からクラウドサブスクリプション収益への転換を進める。"),
        ("顧客基盤拡大", "中堅・中小企業セグメントへ販売チャネルを拡大し、パートナー経由売上比率を40%へ。"),
        ("業務効率化", "全社の文書電子化と承認ワークフロー統合により間接業務時間を20%削減する。"),
    ]
    theme, detail = themes[i % len(themes)]
    if i < 3:
        # docx 版 (RAG 索引対象)
        content = make_docx(
            f"中期経営計画 {year}-{year+2} 重点施策: {theme}",
            [COMPANY, "経営企画 機密文書"],
            [("施策概要", [detail]),
             ("数値目標", [
                 f"売上高: {rnd.randint(80, 120)}億円（{year+2}年度）",
                 f"営業利益率: {rnd.randint(8, 15)}%",
                 f"クラウドサービス売上比率: {rnd.randint(40, 60)}%",
             ]),
             ("リスク", ["競合他社の価格攻勢、採用計画未達、為替変動を主要リスクとして管理する。"])])
        return f"中期経営計画_{year}_{theme}.docx", MIME["docx"], content
    content = make_pptx(
        f"中期経営計画 {year}-{year+2} 説明資料",
        f"{COMPANY} 経営企画",
        [(f"重点施策: {theme}", [detail]),
         ("数値目標", [f"売上高 {rnd.randint(80, 120)}億円", f"営業利益率 {rnd.randint(8, 15)}%"]),
         ("実行体制", ["各本部長をオーナーとするタスクフォースを設置", "四半期毎に取締役会でレビュー"])])
    return f"中期経営計画_{year}_{theme}_説明資料.pptx", MIME["pptx"], content


def gen_proposals(rnd: random.Random, i: int):
    customer = CUSTOMERS[i % len(CUSTOMERS)]
    product = rnd.choice(list(PRODUCTS))
    pain = rnd.choice([
        "紙文書の保管コストと検索性の低さ",
        "契約締結までのリードタイムの長さ",
        "部署毎に分散したファイルサーバの統制不足",
        "監査対応時の証跡収集の負荷",
        "拠点間での文書共有の遅延",
    ])
    effect = rnd.choice([
        f"文書検索時間を{rnd.randint(50, 80)}%削減",
        f"契約締結リードタイムを平均{rnd.randint(5, 14)}日短縮",
        f"保管コストを年間{rnd.randint(2, 8)}百万円削減",
        f"監査準備工数を{rnd.randint(30, 60)}%削減",
    ])
    title = f"{customer}様向け {product}導入提案書"
    if i % 5 == 4:
        content = make_pptx(title, f"{COMPANY} 営業本部 担当: {rnd.choice(SALES_STAFF)}",
                            [("お客様の課題", [pain]),
                             (f"{product}による解決", [PRODUCTS[product], f"期待効果: {effect}"]),
                             ("導入スケジュール", ["要件定義1か月 → 構築2か月 → 並行稼働1か月"]),
                             ("概算費用", [f"初期費用 {rnd.randint(3, 15)}百万円 / 月額 {rnd.randint(10, 80)}万円"])])
        return f"提案書_{customer}_{product}.pptx", MIME["pptx"], content
    content = make_docx(title, [COMPANY, f"担当: {rnd.choice(SALES_STAFF)} / 営業本部"],
                        [("お客様の課題", [f"{customer}様では{pain}が経営課題となっています。"]),
                         (f"ご提案内容: {product}", [
                             f"{PRODUCTS[product]}「{product}」の導入をご提案します。",
                             f"期待効果: {effect}。",
                             f"導入実績: 同業種{rnd.randint(3, 20)}社での稼働実績があります。"]),
                         ("導入スケジュール", ["要件定義（1か月）、環境構築・データ移行（2か月）、並行稼働（1か月）を想定しています。"]),
                         ("概算費用", [f"初期費用{rnd.randint(3, 15)}百万円、月額利用料{rnd.randint(10, 80)}万円（税別）。詳細は見積書をご参照ください。"])])
    return f"提案書_{customer}_{product}.docx", MIME["docx"], content


def gen_quotations(rnd: random.Random, i: int):
    customer = CUSTOMERS[(i * 3 + 1) % len(CUSTOMERS)]
    product = rnd.choice(list(PRODUCTS))
    users = rnd.choice([50, 100, 200, 300, 500])
    unit = rnd.choice([800, 1000, 1200, 1500])
    setup = rnd.randint(2, 10) * 500000
    rows = [
        [1, f"{product} 基本ライセンス（{users}ユーザー・年間）", users, unit * 12, users * unit * 12, "年間サブスクリプション"],
        [2, "初期セットアップ・環境構築", 1, setup, setup, ""],
        [3, "管理者トレーニング", 2, 80000, 160000, "オンライン実施"],
        [4, "既存データ移行支援", rnd.randint(3, 10), 120000, None, "人日単価"],
    ]
    rows[3][4] = rows[3][2] * rows[3][3]
    total = sum(r[4] for r in rows)
    content = make_xlsx(
        "見積書", f"御見積書 EST-2026-{i+1:04d}",
        [f"{customer} 御中", f"件名: {product} 導入一式", f"{COMPANY} 営業本部", "有効期限: 発行日より30日"],
        ["No.", "項目", "数量", "単価（円）", "金額（円）", "備考"],
        rows + [["", "合計（税別）", "", "", total, ""]],
        [f"お支払条件: 検収月末締め翌月末払い。担当: {rnd.choice(SALES_STAFF)}。"])
    return f"見積書_EST-2026-{i+1:04d}_{customer}.xlsx", MIME["xlsx"], content


def gen_visit_reports(rnd: random.Random, i: int):
    customer = CUSTOMERS[(i * 7 + 2) % len(CUSTOMERS)]
    staff = SALES_STAFF[i % len(SALES_STAFF)]
    month = rnd.choice(MONTHS)
    product = rnd.choice(list(PRODUCTS))
    status = rnd.choice(["情報収集", "提案準備", "見積提示済", "クロージング", "既存フォロー"])
    voice = rnd.choice([
        "現行システムの保守期限が迫っており、来期予算での更改を検討中とのこと。",
        "競合製品と比較検討中。価格よりもセキュリティ機能を重視している。",
        "現場部門は乗り気だが情報システム部門の承認プロセスに時間がかかる見込み。",
        "海外拠点への展開も視野に入れており、多言語対応について質問を受けた。",
        "電子帳簿保存法対応が急務であり、年度内の導入を希望している。",
        "稟議は通過済み。役員決裁待ちで、今月中に発注見込み。",
    ])
    next_action = rnd.choice([
        "技術担当同行のうえデモを実施する",
        "セキュリティチェックシートに回答し提出する",
        "見積の再提示（ボリュームディスカウント案）を行う",
        "導入事例資料を送付し、次回訪問日を調整する",
        "PoC環境の提供について社内調整する",
    ])
    content = make_docx(
        f"顧客訪問報告書: {customer}",
        [f"訪問日: {month} / 担当: {staff} / 商談ステータス: {status}", f"対象商材: {product}"],
        [("面談内容", [f"{customer}の情報システム部門および利用部門と面談。{voice}"]),
         ("所感", [f"{product}の{rnd.choice(['検索機能', '権限管理', 'ワークフロー', '監査ログ', 'API連携'])}が評価されており、受注確度は{rnd.choice(['A', 'B', 'C'])}と判断する。"]),
         ("ネクストアクション", [next_action])])
    return f"訪問報告_{month}_{customer}_{staff.split()[0]}.docx", MIME["docx"], content


def gen_sales_reports(rnd: random.Random, i: int):
    month = MONTHS[i % len(MONTHS)]
    section = "東日本営業課" if i % 2 == 0 else "西日本営業課"
    rows = []
    total = 0
    for j in range(6):
        customer = CUSTOMERS[(i * 5 + j) % len(CUSTOMERS)]
        product = list(PRODUCTS)[j % len(PRODUCTS)]
        amount = rnd.randint(80, 600) * 10000
        total += amount
        rows.append([customer, product, amount,
                     rnd.choice(["受注", "受注", "進行中", "提案中", "失注"]),
                     rnd.choice(SALES_STAFF),
                     rnd.choice(["新規", "既存更新", "アップセル"])])
    content = make_xlsx(
        "月次売上", f"{month} {section} 月次売上レポート",
        [COMPANY, f"作成: {section}", f"当月受注合計: {total:,}円"],
        ["顧客名", "商材", "金額（円）", "ステータス", "担当", "区分"],
        rows,
        [f"所感: {rnd.choice(['大型更改案件が寄与し目標を達成', '新規開拓が計画未達のため来月挽回策を実施', 'クラウド移行需要が引き続き旺盛', '失注要因は価格競争力。値引き権限の見直しを提案'])}。"])
    return f"月次売上_{month}_{section}.xlsx", MIME["xlsx"], content


def gen_sales_contracts(rnd: random.Random, i: int):
    customer = CUSTOMERS[(i * 11 + 3) % len(CUSTOMERS)]
    product = rnd.choice(list(PRODUCTS))
    term = rnd.choice(["1年間", "2年間", "3年間"])
    amount = rnd.randint(100, 900) * 10000
    content = make_docx(
        f"サービス利用契約書: {customer}",
        [f"契約番号: CTR-2026-{i+1:04d}", f"{COMPANY}（以下「甲」）と{customer}（以下「乙」）"],
        [("第1条（目的）", [f"甲は乙に対し、{PRODUCTS[product]}「{product}」を提供する。"]),
         ("第2条（契約期間）", [f"契約期間は{term}とし、期間満了の3か月前までに申し出がない場合は同一条件で更新する。"]),
         ("第3条（利用料金）", [f"年額{amount:,}円（税別）。支払は年一括前払いとする。"]),
         ("第4条（秘密保持）", ["両当事者は本契約に関連して知り得た相手方の秘密情報を第三者に開示してはならない。"]),
         ("第5条（サービスレベル）", [f"月間稼働率{rnd.choice(['99.5', '99.9'])}%を保証する。未達の場合は利用料金の減額に応じる。"])])
    return f"契約書_CTR-2026-{i+1:04d}_{customer}.docx", MIME["docx"], content


def gen_design_docs(rnd: random.Random, i: int):
    product = list(PRODUCTS)[i % len(PRODUCTS)]
    component = rnd.choice(["認証基盤", "検索エンジン連携", "帳票出力", "API ゲートウェイ",
                            "監査ログ", "バッチ処理", "通知サービス", "ファイルストレージ"])
    title = f"{product} {component} 設計書 v{i % 3 + 1}.0"
    body, ext = make_pdf_or_docx(
        title,
        [COMPANY + " 技術本部 開発課", f"文書番号: DSN-{i+1:04d}", f"作成: {rnd.choice(DEV_STAFF)}"],
        [("概要", [f"本書は{product}の{component}コンポーネントの設計を記述する。"]),
         ("アーキテクチャ", [
             f"{component}は{rnd.choice(['マイクロサービス', 'モジュラーモノリス'])}構成とし、"
             f"{rnd.choice(['PostgreSQL', 'CouchDB', 'DynamoDB'])}を永続化層に採用する。",
             f"想定スループットは{rnd.randint(100, 2000)}リクエスト/秒、レイテンシ目標はp99で{rnd.randint(100, 800)}ミリ秒。"]),
         ("セキュリティ設計", [
             "通信はTLS1.2以上を必須とし、認可はロールベースアクセス制御で実装する。",
             f"監査ログは{rnd.choice([90, 180, 365])}日間保持する。"]),
         ("障害設計", [f"リトライは指数バックオフ最大{rnd.randint(3, 5)}回。サーキットブレーカで下流障害の伝播を防止する。"])])
    return f"設計書_{product}_{component}_v{i % 3 + 1}.{ext}", MIME[ext], body


def gen_incident_reports(rnd: random.Random, i: int):
    product = list(PRODUCTS)[i % len(PRODUCTS)]
    月 = rnd.choice(MONTHS)
    sev = rnd.choice(["SEV1", "SEV2", "SEV2", "SEV3", "SEV3"])
    cause = rnd.choice([
        ("データベース接続プールの枯渇", "接続リーク箇所を修正し、プール監視アラートを追加した"),
        ("証明書の期限切れ", "証明書自動更新の仕組みを導入し、期限30日前アラートを設定した"),
        ("ストレージ容量逼迫", "ログローテーション設定を見直し、容量監視の閾値を70%に変更した"),
        ("デプロイ後のメモリリーク", "該当リリースをロールバックし、負荷試験項目にヒープ監視を追加した"),
        ("外部API のレート制限超過", "リクエストのバッチ化とバックオフ制御を実装した"),
        ("ネットワーク機器の故障", "冗長経路への自動切替を検証し、フェイルオーバー手順書を更新した"),
    ])
    duration = rnd.randint(10, 300)
    content = make_docx(
        f"障害報告書 INC-2026-{i+1:04d} [{sev}] {product}",
        [f"発生: {月} / 影響時間: {duration}分 / 起票: {rnd.choice(DEV_STAFF)}", COMPANY + " 技術本部"],
        [("事象", [f"{product}で{rnd.choice(['ログイン不可', '検索遅延', 'ファイルアップロード失敗', '帳票出力エラー', '通知遅延'])}が発生し、{rnd.randint(5, 90)}%の顧客に影響した。"]),
         ("原因", [f"直接原因は{cause[0]}。監視の検知が{rnd.randint(3, 30)}分遅れたことが影響拡大の要因となった。"]),
         ("恒久対策", [cause[1] + "。", "再発防止策の完了は次回リリースで確認する。"]),
         ("タイムライン", [f"検知 {rnd.randint(0, 20)}分 → 一次対応開始 → 復旧宣言まで{duration}分。ポストモーテムは翌週実施。"])])
    return f"障害報告_INC-2026-{i+1:04d}_{product}.docx", MIME["docx"], content


def gen_release_notes(rnd: random.Random, i: int):
    product = list(PRODUCTS)[i % len(PRODUCTS)]
    ver = f"{2 + i // 6}.{i % 6}.0"
    feats = rnd.sample([
        "全文検索のインデックス更新を準リアルタイム化",
        "監査ログのCSVエクスポート機能を追加",
        "モバイルアプリのオフライン閲覧に対応",
        "SAML認証のIdP設定ウィザードを追加",
        "帳票テンプレートのバージョン管理に対応",
        "APIレート制限の管理画面を追加",
        "ダッシュボードのウィジェットを刷新",
        "ファイルプレビューの対応形式を拡充",
    ], 3)
    fixes = rnd.sample([
        "大容量ファイルアップロード時のタイムアウトを修正",
        "特定条件で検索結果の並び順が不定になる問題を修正",
        "通知メールの文字化けを修正",
        "権限変更が即時反映されない問題を修正",
    ], 2)
    title = f"{product} v{ver} リリースノート"
    body, ext = make_pdf_or_docx(
        title, [COMPANY + " 技術本部", f"リリース日: {rnd.choice(MONTHS)}"],
        [("新機能", feats), ("不具合修正", fixes),
         ("アップグレード時の注意", [f"データベーススキーマの更新を伴うため、事前バックアップを必須とする。想定停止時間は{rnd.choice([10, 30, 60])}分。"])])
    return f"リリースノート_{product}_v{ver}.{ext}", MIME[ext], body


def gen_runbooks(rnd: random.Random, i: int):
    target = rnd.choice(["本番データベース", "検索クラスタ", "アプリケーションサーバ", "ロードバランサ",
                         "バックアップ基盤", "監視基盤"])
    op = rnd.choice(["定期再起動", "フェイルオーバー切替", "バックアップリストア", "証明書更新",
                     "スケールアウト", "パッチ適用"])
    title = f"運用手順書: {target} {op}"
    content = make_docx(
        title, [COMPANY + " 技術本部 インフラ課", f"手順書番号: OPS-{i+1:03d}", "承認: 平田技術本部長"],
        [("前提条件", [f"作業は{rnd.choice(['メンテナンスウィンドウ（土曜22:00-翌6:00）', '平日夜間（21:00以降）'])}に実施する。",
                   "作業前にオンコール担当へ連絡し、監視の一時抑止を設定する。"]),
         ("手順", [
             f"1. {target}の事前ヘルスチェックを実施し、結果を作業チケットに記録する。",
             f"2. {op}を実施する。想定所要時間は{rnd.randint(15, 120)}分。",
             "3. 事後ヘルスチェックとスモークテストを実施する。",
             "4. 監視抑止を解除し、作業完了を報告する。"]),
         ("ロールバック", [f"手順2で異常が発生した場合は直前のスナップショットから復旧する。復旧目標時間は{rnd.randint(30, 90)}分。"])])
    return f"運用手順_{i+1:03d}_{target}_{op}.docx", MIME["docx"], content


def gen_tech_research(rnd: random.Random, i: int):
    topic = rnd.choice([
        ("ベクトル検索エンジンの比較評価", "埋め込みモデルの次元数と検索精度・レイテンシのトレードオフを評価した。"),
        ("大規模言語モデルの社内活用", "文書要約と問い合わせ応答への適用可能性をPoCで検証した。"),
        ("コンテナオーケストレーションの移行", "現行構成からKubernetesへの移行コストと運用負荷を試算した。"),
        ("ゼロトラストネットワーク", "境界型防御からの移行ステップとIDaaS製品の比較を行った。"),
        ("サーバレスアーキテクチャ", "バッチ処理のサーバレス化によるコスト削減効果を測定した。"),
        ("OCR・帳票認識エンジン", "手書き帳票の認識精度を3製品で比較評価した。"),
        ("データレイクハウス", "分析基盤の統合に向けたストレージフォーマットの選定を行った。"),
    ])
    title = f"技術調査報告: {topic[0]}"
    content = make_docx(
        title, [COMPANY + " 技術本部", f"調査担当: {rnd.choice(DEV_STAFF)}", f"実施: {rnd.choice(MONTHS)}"],
        [("調査目的", [topic[1]]),
         ("結果", [
             f"評価スコアは候補Aが{rnd.randint(70, 95)}点、候補Bが{rnd.randint(60, 90)}点、候補Cが{rnd.randint(50, 85)}点となった。",
             f"総所有コストは3年間で最大{rnd.randint(20, 60)}%の差が生じる試算となった。"]),
         ("推奨", [rnd.choice(["候補Aを本採用としPoCを次四半期に開始することを推奨する。",
                            "現時点では時期尚早と判断し、半年後に再評価する。",
                            "限定的なワークロードから段階導入することを推奨する。"])])])
    return f"技術調査_{i+1:02d}_{topic[0][:12]}.docx", MIME["docx"], content


def gen_hr_evaluations(rnd: random.Random, i: int):
    period = rnd.choice(["2025年下期", "2026年上期"])
    theme = rnd.choice([
        ("評価制度改定案", "コンピテンシー評価と成果評価の比率を50:50から40:60に変更する案を検討する。"),
        ("昇格候補者リスト", "課長昇格の候補者について、評価履歴と360度フィードバックの結果を整理した。"),
        ("評価者研修計画", "評価の甘辛差を是正するため、全評価者に対しキャリブレーション研修を実施する。"),
        ("目標設定ガイドライン", "OKR形式の目標設定を営業本部と技術本部で先行導入する。"),
        ("エンゲージメント調査結果", f"全社スコアは{rnd.randint(60, 78)}ポイント。上司との1on1頻度と相関が見られた。"),
        ("評価分布レポート", f"S評価{rnd.randint(3, 8)}%、A評価{rnd.randint(20, 30)}%、B評価{rnd.randint(45, 60)}%の分布となった。"),
    ])
    title = f"人事評価: {theme[0]}（{period}）"
    content = make_docx(
        title, [COMPANY + " 人事課 取扱注意", "本文書は人事課および経営会議限りとする。"],
        [("内容", [theme[1]]),
         ("留意事項", ["個人が特定される情報の取り扱いは人事課内に限定する。",
                    f"次回見直しは{rnd.choice(MONTHS)}を予定。"])])
    return f"人事評価_{period}_{theme[0]}.docx", MIME["docx"], content


def gen_hr_salary(rnd: random.Random, i: int):
    period = rnd.choice(["2026年4月改定", "2026年賞与"])
    grades = ["G1(担当)", "G2(主任)", "G3(課長)", "G4(本部長)"]
    rows = []
    for g in grades:
        base = {"G1(担当)": 280000, "G2(主任)": 350000, "G3(課長)": 480000, "G4(本部長)": 650000}[g]
        rows.append([g, base, int(base * (1 + rnd.randint(2, 5) / 100)),
                     f"{rnd.randint(2, 5)}%", rnd.choice(["物価上昇分を反映", "市場水準に合わせ引上げ", "据え置きから改定"])])
    content = make_xlsx(
        "給与改定", f"給与テーブル改定案（{period}） 第{i+1}版",
        [COMPANY + " 人事課 【厳秘】", "本資料は人事課外への持ち出しを禁止する。賞与原資および昇給原資の試算を含む。"],
        ["等級", "現行基本給（円）", "改定後基本給（円）", "改定率", "備考"],
        rows,
        [f"賞与原資: 基本給総額の{rnd.choice([4.2, 4.5, 5.0])}か月分で試算。",
         f"昇給原資総額: {rnd.randint(18, 40)}百万円。取締役会承認後に適用する。"])
    return f"給与テーブル改定_{period}_第{i+1}版.xlsx", MIME["xlsx"], content


def gen_invoices(rnd: random.Random, i: int):
    customer = CUSTOMERS[(i * 13 + 5) % len(CUSTOMERS)]
    product = rnd.choice(list(PRODUCTS))
    amount = rnd.randint(30, 500) * 10000
    tax = amount // 10
    month = MONTHS[i % len(MONTHS)]
    content = make_xlsx(
        "請求書", f"請求書 INV-2026-{i+1:04d}",
        [f"{customer} 御中", f"請求月: {month} / 支払期限: 翌月末", COMPANY + " 経理課"],
        ["No.", "品目", "金額（円）", "備考"],
        [[1, f"{product} 利用料（{month}分）", amount, PRODUCTS[product]],
         [2, "消費税（10%）", tax, ""],
         ["", "合計", amount + tax, ""]],
        ["振込先: 北都銀行 本店 普通 1234567 ヒナタサンギョウ（カ"])
    return f"請求書_INV-2026-{i+1:04d}_{customer}.xlsx", MIME["xlsx"], content


def gen_budgets(rnd: random.Random, i: int):
    dept = ["営業本部", "技術本部", "管理本部", "全社"][i % 4]
    year = "2026年度"
    items = {
        "営業本部": [("人件費", 180), ("販売促進費", 45), ("出張旅費", 30), ("展示会出展", 20)],
        "技術本部": [("人件費", 260), ("クラウド利用料", 80), ("開発機材", 25), ("外部委託", 60)],
        "管理本部": [("人件費", 90), ("オフィス賃借料", 70), ("採用費", 25), ("システム利用料", 15)],
        "全社": [("人件費合計", 530), ("設備投資", 120), ("研究開発", 95), ("予備費", 40)],
    }[dept]
    rows = []
    for name, base in items:
        budget = base * 1000000
        actual = int(budget * rnd.uniform(0.7, 1.1))
        rows.append([name, budget, actual, f"{actual / budget * 100:.0f}%",
                     rnd.choice(["計画どおり", "上振れ注意", "下振れ（採用遅延）", "期ずれあり"])])
    content = make_xlsx(
        "予算", f"{year} {dept} 予算執行状況（第{i // 4 + 1}四半期）",
        [COMPANY + " 経理課", "経営会議報告用"],
        ["費目", "年間予算（円）", "執行額（円）", "執行率", "コメント"],
        rows,
        ["執行率90%超の費目は翌四半期の支出計画見直しを要する。"])
    return f"予算執行_{year}_{dept}_Q{i // 4 + 1}.xlsx", MIME["xlsx"], content


def gen_legal_reviews(rnd: random.Random, i: int):
    customer = CUSTOMERS[(i * 17 + 7) % len(CUSTOMERS)]
    kind = rnd.choice(["サービス利用契約", "業務委託契約", "秘密保持契約", "販売代理店契約", "共同開発契約"])
    risk = rnd.choice([
        ("損害賠償の上限規定がない", "賠償上限を直近12か月の受領額とする条項の追加を求める"),
        ("知的財産権の帰属が不明確", "成果物の著作権帰属と利用許諾範囲を明記するよう修正する"),
        ("解約時のデータ返還義務が過大", "返還期限を30日から90日に緩和する対案を提示する"),
        ("準拠法が相手方所在地法", "準拠法を日本法、専属的合意管轄を東京地方裁判所とするよう求める"),
        ("再委託の事前承諾条項が営業実態と不整合", "包括承諾方式への変更を提案する"),
    ])
    content = make_docx(
        f"契約審査報告: {customer} {kind}",
        [COMPANY + " 法務課", f"審査番号: LGL-2026-{i+1:04d}", f"審査担当: 馬場 遼"],
        [("審査対象", [f"{customer}との{kind}（先方雛形）について審査を実施した。"]),
         ("指摘事項", [f"第{rnd.randint(5, 20)}条: {risk[0]}。リスク評価は{rnd.choice(['高', '中'])}。"]),
         ("修正方針", [risk[1] + "。営業担当は先方との交渉時に本方針に従うこと。"]),
         ("結論", [rnd.choice(["条件付き承認（上記修正を必須とする）", "再審査要（修正案の再提出を求める）", "承認（軽微な指摘のみ）"])])])
    return f"契約審査_LGL-2026-{i+1:04d}_{customer}.docx", MIME["docx"], content


def gen_projx_specs(rnd: random.Random, i: int):
    feature = rnd.choice(["自然言語検索", "回答生成", "権限連動フィルタ", "多言語対応",
                          "既存製品連携", "利用状況分析", "オンプレミス版", "監査対応"])
    content = make_docx(
        f"新製品X({PROJX_CODENAME}) 製品仕様書: {feature}",
        [COMPANY + f" 新製品Xプロジェクト 【極秘】", "本文書はプロジェクトメンバー以外への開示を禁止する。"],
        [("製品コンセプト", [f"{PROJX_CODENAME}は{PROJX_DESC}であり、社内文書に対する自然言語での質問応答を提供する。"]),
         (f"機能仕様: {feature}", [
             f"{feature}は{rnd.choice(['初回リリース', 'v1.1', 'v2.0'])}のスコープとする。",
             f"ベクトル検索と権限フィルタを組み合わせ、利用者の閲覧権限内でのみ回答を生成する。",
             f"応答時間目標は{rnd.choice([2, 3, 5])}秒以内（p95）。"]),
         ("既存製品との関係", [f"DocuHiveの文書ストアを検索対象とし、SignFlowの契約書検索にも拡張予定。"])])
    return f"AuroraSpec_{i+1:02d}_{feature}.docx", MIME["docx"], content


def gen_projx_pricing(rnd: random.Random, i: int):
    plan = ["スターター", "ビジネス", "エンタープライズ"]
    rows = []
    for j, p in enumerate(plan):
        monthly = [50000, 180000, 500000][j]
        monthly = int(monthly * rnd.uniform(0.9, 1.15))
        rows.append([p, monthly, [100, 500, 5000][j],
                     ["基本検索のみ", "回答生成・API連携", "専有環境・SLA99.9%"][j]])
    content = make_xlsx(
        "価格戦略", f"新製品X({PROJX_CODENAME}) 価格シミュレーション 案{i+1}",
        [COMPANY + " 新製品Xプロジェクト 【極秘】",
         f"前提: 初年度{rnd.randint(30, 120)}社獲得、解約率月{rnd.uniform(0.5, 2.0):.1f}%"],
        ["プラン", "月額（円）", "対象ユーザー数", "含まれる機能"],
        rows,
        [f"競合の同等プラン比で{rnd.randint(5, 25)}%低い価格設定とし、既存DocuHive顧客への同時提案で獲得コストを抑える。",
         "正式価格は取締役会承認後に確定する。本シートの社外共有を禁止する。"])
    return f"Aurora価格シミュレーション_案{i+1}.xlsx", MIME["xlsx"], content


def gen_projx_market(rnd: random.Random, i: int):
    seg = rnd.choice(["金融", "製造", "自治体", "医療", "教育"])
    if i % 3 == 2:
        content = make_pptx(
            f"新製品X({PROJX_CODENAME}) 市場調査: {seg}セグメント",
            f"{COMPANY} 新製品Xプロジェクト 【極秘】",
            [("市場規模", [f"{seg}向けAI文書検索市場は2028年に{rnd.randint(200, 800)}億円規模と推計"]),
             ("競合", ["先行2社は精度面で優位だが権限連動に弱み", f"{PROJX_CODENAME}は既存文書管理基盤との統合を差別化軸とする"]),
             ("参入戦略", [f"{seg}の既存顧客{rnd.randint(5, 30)}社へのクローズドベータ提供から開始"])])
        return f"Aurora市場調査_{seg}.pptx", MIME["pptx"], content
    content = make_docx(
        f"新製品X({PROJX_CODENAME}) 市場調査報告: {seg}セグメント",
        [COMPANY + " 新製品Xプロジェクト 【極秘】", f"調査実施: {rnd.choice(MONTHS)}"],
        [("市場規模", [f"{seg}セグメントにおけるAI文書検索・質問応答市場は2028年に{rnd.randint(200, 800)}億円規模に達すると推計される。"]),
         ("競合分析", [f"先行競合は{rnd.randint(2, 4)}社。検索精度では先行するが、文書権限と連動した回答制御は未対応であり、{PROJX_CODENAME}の差別化要因となる。"]),
         ("想定顧客の声", [f"{seg}業界のヒアリングでは「{rnd.choice(['監査対応の証跡検索', '過去の稟議・契約の横断検索', '規程類の問い合わせ削減', 'ナレッジ継承'])}」への期待が最も高かった。"])])
    return f"Aurora市場調査_{i+1:02d}_{seg}.docx", MIME["docx"], content


# ---------------------------------------------------------------------------
# カテゴリ → (件数, ジェネレータ)。合計 300 件。
# ---------------------------------------------------------------------------
CATEGORY_SPECS: dict[str, tuple[int, Callable[[random.Random, int], tuple[str, str, bytes]]]] = {
    "rules": (14, gen_rules),
    "announcements": (12, gen_announcements),
    "board_minutes": (12, gen_board_minutes),
    "mid_term_plan": (6, gen_mid_term_plan),
    "proposals": (18, gen_proposals),
    "quotations": (18, gen_quotations),
    "visit_reports": (24, gen_visit_reports),
    "sales_reports": (16, gen_sales_reports),
    "sales_contracts": (12, gen_sales_contracts),
    "design_docs": (16, gen_design_docs),
    "incident_reports": (20, gen_incident_reports),
    "release_notes": (14, gen_release_notes),
    "runbooks": (12, gen_runbooks),
    "tech_research": (14, gen_tech_research),
    "hr_evaluations": (14, gen_hr_evaluations),
    "hr_salary": (8, gen_hr_salary),
    "invoices": (20, gen_invoices),
    "budgets": (10, gen_budgets),
    "legal_reviews": (14, gen_legal_reviews),
    "projx_specs": (10, gen_projx_specs),
    "projx_pricing": (8, gen_projx_pricing),
    "projx_market": (8, gen_projx_market),
}

TOTAL_DOCS = sum(n for n, _ in CATEGORY_SPECS.values())
assert TOTAL_DOCS == 300, f"category counts must sum to 300, got {TOTAL_DOCS}"


def iter_documents(categories: list[str] | None = None) -> Iterator[GeneratedDoc]:
    """全カテゴリ (または指定カテゴリ) の文書を生成して yield する。"""
    for category, (count, gen) in CATEGORY_SPECS.items():
        if categories and category not in categories:
            continue
        rnd = random.Random(f"{category}:{SEED}")
        seen: set[str] = set()
        for i in range(count):
            filename, mimetype, content = gen(rnd, i)
            # 万一の名前衝突は連番サフィックスで回避 (同一フォルダ内一意制約)
            if filename in seen:
                stem, dot, ext = filename.rpartition(".")
                filename = f"{stem}_{i+1:02d}{dot}{ext}"
            seen.add(filename)
            yield GeneratedDoc(category, filename, mimetype, content)


if __name__ == "__main__":
    # スモークテスト: 全件生成してカテゴリ毎の件数とサイズを表示
    from collections import Counter

    counter: Counter = Counter()
    size = 0
    for d in iter_documents():
        counter[d.category] += 1
        size += len(d.content)
    for cat, n in counter.items():
        print(f"{cat:>18}: {n}")
    print(f"total: {sum(counter.values())} docs, {size / 1024 / 1024:.1f} MB")
