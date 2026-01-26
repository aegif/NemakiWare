#!/usr/bin/env python3
"""
NemakiWare Test Environment Setup Script

This script populates a clean NemakiWare instance with test data including:
- Organizations (groups)
- Users
- Folder structures with permissions
- Custom content types
- Test documents (Word, PowerPoint, PDF, Excel)

Usage:
    python setup_test_environment.py [--config config.yaml] [--seed 12345]

Environment Variables (override config file values):
    NEMAKI_BROWSER_URL: CMIS Browser Binding URL (e.g., http://localhost:8080/core/browser/bedroom)
    NEMAKI_CMIS_USERNAME: CMIS username
    NEMAKI_CMIS_PASSWORD: CMIS password
    NEMAKI_REST_URL: REST API base URL
    NEMAKI_REST_USERNAME: REST API username
    NEMAKI_REST_PASSWORD: REST API password
    NEMAKI_RANDOM_SEED: Random seed for reproducible test data generation
"""

import argparse
import io
import mimetypes
import os
import random
import sys
from datetime import datetime, timedelta

import requests
import yaml
from requests.auth import HTTPBasicAuth
from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

# HTTP request timeout settings (connect timeout, read timeout) in seconds
HTTP_TIMEOUT = (10, 30)


class TestEnvironmentSetup:
    """Main class for setting up the NemakiWare test environment."""

    # Japanese font name for PDF generation (registered once at initialization)
    JAPANESE_FONT = "HeiseiMin-W3"

    def __init__(self, config_path: str):
        """Initialize with configuration file path."""
        self.config = self._load_config(config_path)
        self.root_folder_id = None
        self.sites_folder_id = None
        self.created_groups = {}
        self.created_users = {}
        self.org_folder_ids = {}

        # Register Japanese CID font once at initialization (built-in, no external font file needed)
        pdfmetrics.registerFont(UnicodeCIDFont(self.JAPANESE_FONT))

    def _load_config(self, config_path: str) -> dict:
        """Load configuration from YAML file and apply environment variable overrides."""
        with open(config_path, "r", encoding="utf-8") as f:
            config = yaml.safe_load(f)

        # Apply environment variable overrides for credentials
        if os.environ.get("NEMAKI_BROWSER_URL"):
            config["cmis"]["browser_url"] = os.environ["NEMAKI_BROWSER_URL"]
        if os.environ.get("NEMAKI_CMIS_USERNAME"):
            config["cmis"]["username"] = os.environ["NEMAKI_CMIS_USERNAME"]
        if os.environ.get("NEMAKI_CMIS_PASSWORD"):
            config["cmis"]["password"] = os.environ["NEMAKI_CMIS_PASSWORD"]
        if os.environ.get("NEMAKI_REST_URL"):
            config["rest_api"]["base_url"] = os.environ["NEMAKI_REST_URL"]
        if os.environ.get("NEMAKI_REST_USERNAME"):
            config["rest_api"]["username"] = os.environ["NEMAKI_REST_USERNAME"]
        if os.environ.get("NEMAKI_REST_PASSWORD"):
            config["rest_api"]["password"] = os.environ["NEMAKI_REST_PASSWORD"]
        if os.environ.get("NEMAKI_REPOSITORY_ID"):
            config["cmis"]["repository_id"] = os.environ["NEMAKI_REPOSITORY_ID"]

        return config

    def _get_cmis_auth(self) -> HTTPBasicAuth:
        """Get CMIS authentication."""
        cmis_config = self.config["cmis"]
        return HTTPBasicAuth(cmis_config["username"], cmis_config["password"])

    def _get_browser_url(self) -> str:
        """Get CMIS Browser Binding URL."""
        return self.config["cmis"]["browser_url"].rstrip("/")

    def _get_rest_auth(self) -> tuple:
        """Get REST API authentication tuple."""
        rest_config = self.config["rest_api"]
        return (rest_config["username"], rest_config["password"])

    def _get_rest_url(self, path: str) -> str:
        """Get full REST API URL."""
        base_url = self.config["rest_api"]["base_url"].rstrip("/")
        repo_id = self.config["cmis"]["repository_id"]
        return f"{base_url}/repo/{repo_id}/{path}"

    def connect_cmis(self) -> None:
        """Connect to CMIS repository using Browser Binding."""
        browser_url = self._get_browser_url()
        print(f"Connecting to CMIS repository at {browser_url}...")

        # Get repository info
        response = requests.get(
            f"{browser_url}?cmisselector=repositoryInfo",
            auth=self._get_cmis_auth(),
            timeout=HTTP_TIMEOUT
        )
        response.raise_for_status()

        repo_info = response.json()
        repo_id = self.config["cmis"]["repository_id"]

        if repo_id in repo_info:
            self.root_folder_id = repo_info[repo_id]["rootFolderId"]
            print(f"Connected to repository: {repo_id}")
            print(f"Root folder ID: {self.root_folder_id}")
        else:
            raise ValueError(f"Repository '{repo_id}' not found")

    def _get_folder_children(self, folder_id: str) -> list:
        """Get children of a folder using CMIS 1.1 Browser Binding."""
        browser_url = self._get_browser_url()
        response = requests.get(
            f"{browser_url}/root?cmisselector=children&objectId={folder_id}",
            auth=self._get_cmis_auth(),
            timeout=HTTP_TIMEOUT
        )
        if response.status_code == 200:
            return response.json().get("objects", [])
        return []

    def _find_folder_by_name(self, parent_id: str, name: str) -> str:
        """Find a folder by name in parent. Returns folder ID or None."""
        for obj in self._get_folder_children(parent_id):
            props = obj["object"]["properties"]
            if props["cmis:name"]["value"] == name:
                return props["cmis:objectId"]["value"]
        return None

    def _create_folder(self, parent_id: str, name: str) -> str:
        """Create a folder. Returns folder ID."""
        # Check if folder already exists
        existing = self._find_folder_by_name(parent_id, name)
        if existing:
            print(f"  Found existing folder: {name}")
            return existing

        browser_url = self._get_browser_url()
        response = requests.post(
            browser_url,
            auth=self._get_cmis_auth(),
            data={
                "cmisaction": "createFolder",
                "folderId": parent_id,
                "propertyId[0]": "cmis:objectTypeId",
                "propertyValue[0]": "cmis:folder",
                "propertyId[1]": "cmis:name",
                "propertyValue[1]": name,
            },
            timeout=HTTP_TIMEOUT
        )

        # Accept both 200 and 201 status codes for successful creation
        if response.status_code in (200, 201):
            try:
                result = response.json()
                if "succinctProperties" in result:
                    folder_id = result["succinctProperties"]["cmis:objectId"]
                    print(f"  Created folder: {name}")
                    return folder_id
                elif "properties" in result:
                    folder_id = result["properties"]["cmis:objectId"]["value"]
                    print(f"  Created folder: {name}")
                    return folder_id
            except Exception:
                pass

        # Check for already exists error
        if "contentAlreadyExistsException" in response.text:
            existing = self._find_folder_by_name(parent_id, name)
            if existing:
                print(f"  Found existing folder: {name}")
                return existing

        # Try to extract folder ID from response even on non-200 status
        try:
            result = response.json()
            if "succinctProperties" in result:
                folder_id = result["succinctProperties"]["cmis:objectId"]
                print(f"  Created folder: {name}")
                return folder_id
        except Exception:
            pass

        print(f"  Error creating folder {name}: {response.text[:200]}")
        return None

    def _find_or_create_folder(self, parent_id: str, folder_name: str) -> str:
        """Find existing folder or create new one."""
        return self._create_folder(parent_id, folder_name)

    def _get_sites_folder(self) -> str:
        """Get or create the Sites folder."""
        if self.sites_folder_id is None:
            # NemakiWare's default Sites folder uses capital S
            self.sites_folder_id = self._find_or_create_folder(self.root_folder_id, "Sites")
        return self.sites_folder_id

    def create_groups(self) -> None:
        """Create organization groups."""
        print("\n=== Creating Groups ===")

        for org in self.config["organizations"]:
            group_id = org["id"]
            group_name = org["name"]

            url = self._get_rest_url(f"group/create/{group_id}")
            data = {"name": group_name, "users": "[]", "groups": "[]"}

            try:
                response = requests.post(url, data=data, auth=self._get_rest_auth(), timeout=HTTP_TIMEOUT)
                response.raise_for_status()
                result = response.json()

                if result.get("status") == "success":
                    print(f"  Created group: {group_name} ({group_id})")
                    self.created_groups[group_id] = group_name
                else:
                    error = result.get("error", "Unknown error")
                    if "alreadyExists" in str(error):
                        print(f"  Group already exists: {group_name} ({group_id})")
                        self.created_groups[group_id] = group_name
                    else:
                        print(f"  Failed to create group {group_name}: {error}")
            except requests.exceptions.RequestException as e:
                print(f"  Error creating group {group_name}: {e}")

        self._setup_group_hierarchy()

    def _setup_group_hierarchy(self) -> None:
        """Set up parent-child relationships between groups."""
        print("\n=== Setting up Group Hierarchy ===")

        for org in self.config["organizations"]:
            if org.get("parent"):
                child_id = org["id"]
                parent_id = org["parent"]

                url = self._get_rest_url(f"group/add/{parent_id}")
                data = {"users": "[]", "groups": f'["{child_id}"]'}

                try:
                    response = requests.put(url, data=data, auth=self._get_rest_auth(), timeout=HTTP_TIMEOUT)
                    response.raise_for_status()
                    result = response.json()

                    if result.get("status") == "success":
                        print(f"  Added {child_id} to {parent_id}")
                    else:
                        error = result.get("error", "Unknown error")
                        if "alreadyMember" not in str(error):
                            print(f"  Failed to add {child_id} to {parent_id}: {error}")
                except requests.exceptions.RequestException as e:
                    print(f"  Error setting up hierarchy: {e}")

    def create_users(self) -> None:
        """Create users for each organization."""
        print("\n=== Creating Users ===")

        users_per_org = self.config["users"]["per_organization"]
        password = self.config["users"]["password"]

        for org in self.config["organizations"]:
            org_id = org["id"]
            org_name = org["name"]

            for i in range(1, users_per_org + 1):
                user_id = f"{org_id}{i}"
                user_name = f"{org_name}{i}"

                url = self._get_rest_url(f"user/create/{user_id}")
                data = {
                    "name": user_name,
                    "password": password,
                    "firstName": org_name,
                    "lastName": str(i),
                    "email": f"{user_id}@example.com",
                }

                try:
                    response = requests.post(url, data=data, auth=self._get_rest_auth(), timeout=HTTP_TIMEOUT)
                    response.raise_for_status()
                    result = response.json()

                    if result.get("status") == "success":
                        print(f"  Created user: {user_name} ({user_id})")
                        self.created_users[user_id] = user_name
                    else:
                        error = result.get("error", "Unknown error")
                        if "alreadyExists" in str(error):
                            print(f"  User already exists: {user_name} ({user_id})")
                            self.created_users[user_id] = user_name
                        else:
                            print(f"  Failed to create user {user_name}: {error}")
                except requests.exceptions.RequestException as e:
                    print(f"  Error creating user {user_name}: {e}")

                self._add_user_to_group(user_id, org_id)

    def _add_user_to_group(self, user_id: str, group_id: str) -> None:
        """Add a user to a group."""
        url = self._get_rest_url(f"group/add/{group_id}")
        data = {"users": f'["{user_id}"]', "groups": "[]"}

        try:
            response = requests.put(url, data=data, auth=self._get_rest_auth(), timeout=HTTP_TIMEOUT)
            response.raise_for_status()
            result = response.json()

            if result.get("status") != "success":
                error = result.get("error", "Unknown error")
                if "alreadyMember" not in str(error):
                    print(f"    Failed to add {user_id} to {group_id}: {error}")
        except requests.exceptions.RequestException as e:
            print(f"    Error adding user to group: {e}")

    def register_custom_types(self) -> None:
        """Register custom content types."""
        print("\n=== Registering Custom Types ===")

        for type_def in self.config["custom_types"]:
            self._register_type(type_def)

    def _register_type(self, type_def: dict) -> None:
        """Register a single custom type."""
        type_id = type_def["type_id"]
        print(f"  Registering type: {type_id}")

        xml_content = self._generate_type_xml(type_def)

        url = self._get_rest_url("type/register")

        try:
            files = {"data": ("type.xml", xml_content, "application/xml")}
            response = requests.post(url, files=files, auth=self._get_rest_auth(), timeout=HTTP_TIMEOUT)
            response.raise_for_status()
            result = response.json()

            if result.get("status") == "success":
                print(f"    Type registered successfully: {type_id}")
            else:
                error = result.get("error", "Unknown error")
                if "alreadyExists" in str(error) or "already exists" in str(error):
                    print(f"    Type already exists: {type_id}")
                else:
                    print(f"    Failed to register type {type_id}: {error}")
        except requests.exceptions.RequestException as e:
            print(f"    Error registering type {type_id}: {e}")

    def _generate_type_xml(self, type_def: dict) -> str:
        """Generate XML for type registration."""
        type_id = type_def["type_id"]
        parent_type = type_def["parent_type"]

        if parent_type == "nemaki:document":
            parent_xml = "cm:content"
        elif parent_type == "cmis:folder":
            parent_xml = "cm:folder"
        else:
            parent_xml = parent_type

        properties_xml = ""
        for prop in type_def.get("properties", []):
            prop_type = prop["type"]
            if prop_type == "string":
                data_type = "d:text"
            elif prop_type == "integer":
                data_type = "d:long"
            elif prop_type == "boolean":
                data_type = "d:boolean"
            elif prop_type == "datetime":
                data_type = "d:datetime"
            elif prop_type == "decimal":
                data_type = "d:double"
            else:
                data_type = "d:text"

            properties_xml += f"""
            <property name="{prop['id']}">
                <title>{prop['name']}</title>
                <type>{data_type}</type>
                <mandatory>{"true" if prop.get("required", False) else "false"}</mandatory>
                <index enabled="{"true" if prop.get("queryable", True) else "false"}"/>
            </property>"""

        xml = f"""<?xml version="1.0" encoding="UTF-8"?>
<model>
    <types>
        <type name="{type_id}">
            <title>{type_def['display_name']}</title>
            <parent>{parent_xml}</parent>
            <properties>{properties_xml}
            </properties>
        </type>
    </types>
</model>"""
        return xml

    def create_organization_folders(self) -> None:
        """Create organization folders under /sites."""
        print("\n=== Creating Organization Folders ===")

        if not self.config["sites_folders"]["create_org_folders"]:
            print("  Skipping organization folders (disabled in config)")
            return

        sites_folder_id = self._get_sites_folder()

        for org in self.config["organizations"]:
            if org.get("parent") is None:
                self._create_org_folder_recursive(sites_folder_id, org)

    def _create_org_folder_recursive(self, parent_id: str, org: dict) -> None:
        """Recursively create organization folder and its children."""
        org_id = org["id"]
        org_name = org["name"]

        folder_id = self._find_or_create_folder(parent_id, org_name)
        if folder_id:
            self.org_folder_ids[org_id] = folder_id
            self._set_folder_permissions_for_org(folder_id, org_id)

            for child_org in self.config["organizations"]:
                if child_org.get("parent") == org_id:
                    self._create_org_folder_recursive(folder_id, child_org)

    def _set_folder_permissions_for_org(self, folder_id: str, org_id: str) -> None:
        """Set folder permissions for organization members using Browser Binding."""
        try:
            browser_url = self._get_browser_url()

            # Apply ACL using Browser Binding
            response = requests.post(
                browser_url,
                auth=self._get_cmis_auth(),
                data={
                    "cmisaction": "applyACL",
                    "objectId": folder_id,
                    "ACLPropagation": "propagate",
                    "addACEPrincipal[0]": org_id,
                    "addACEPermission[0][0]": "cmis:read",
                    "addACEPrincipal[1]": org_id,
                    "addACEPermission[1][0]": "cmis:write",
                },
                timeout=HTTP_TIMEOUT
            )

            if response.status_code == 200:
                print(f"    Set permissions for {org_id} on folder")
            else:
                print(f"    Warning: Could not set permissions for {org_id}")
        except Exception as e:
            print(f"    Error setting permissions: {e}")

    def create_root_folders(self) -> None:
        """Create root level folders with content."""
        print("\n=== Creating Root Level Folders ===")

        for folder_config in self.config["root_folders"]:
            self._create_root_folder_with_content(folder_config)

    def _create_root_folder_with_content(self, folder_config: dict) -> None:
        """Create a root folder and populate it with content."""
        folder_name = folder_config["name"]
        print(f"\n  Creating folder: {folder_name}")

        folder_id = self._find_or_create_folder(self.root_folder_id, folder_name)
        if not folder_id:
            print(f"    Failed to create folder: {folder_name}")
            return

        self._set_folder_permissions(folder_id, folder_config.get("permissions", []))

        content_config = folder_config.get("content")
        if content_config:
            self._populate_folder_with_content(folder_id, folder_config)

    def _set_folder_permissions(self, folder_id: str, permissions: list) -> None:
        """Set permissions on a folder using Browser Binding."""
        if not permissions:
            return

        try:
            browser_url = self._get_browser_url()

            # Build ACL data
            data = {
                "cmisaction": "applyACL",
                "objectId": folder_id,
                "ACLPropagation": "propagate",
            }

            for i, perm in enumerate(permissions):
                principal = perm["principal"]
                permission = perm["permission"]
                data[f"addACEPrincipal[{i}]"] = principal
                data[f"addACEPermission[{i}][0]"] = permission

            response = requests.post(
                browser_url,
                auth=self._get_cmis_auth(),
                data=data,
                timeout=HTTP_TIMEOUT
            )

            if response.status_code == 200:
                print(f"    Set {len(permissions)} permission entries")
            else:
                print(f"    Warning: Could not set permissions")
        except Exception as e:
            print(f"    Error setting permissions: {e}")

    def _populate_folder_with_content(self, folder_id: str, folder_config: dict) -> None:
        """Populate a folder with generated content."""
        content_config = folder_config["content"]
        folder_conf_id = folder_config["id"]
        content_type = content_config["type"]
        count = content_config["count"]
        formats = content_config["formats"]
        pages = content_config.get("pages", 1)

        print(f"    Generating {count} documents of type {content_type}...")

        for i in range(count):
            format_choice = self._choose_format(formats)
            extension = format_choice["extension"]

            if folder_conf_id == "shanai_kitei":
                doc_data = self._generate_internal_regulation(i, extension, pages)
            elif folder_conf_id == "keiyakusho":
                # Contracts are always PDF format
                doc_data = self._generate_contract(i, pages)
            elif folder_conf_id == "invoice":
                doc_data = self._generate_invoice(i, extension)
            else:
                continue

            self._upload_document(
                folder_id, doc_data["name"], doc_data["content"], content_type, doc_data.get("properties", {})
            )

            if (i + 1) % 10 == 0:
                print(f"      Created {i + 1}/{count} documents")

        print(f"    Completed: {count} documents created")

    def _choose_format(self, formats: list) -> dict:
        """Choose a format based on ratio weights."""
        total = sum(f["ratio"] for f in formats)
        r = random.random() * total
        cumulative = 0
        for f in formats:
            cumulative += f["ratio"]
            if r <= cumulative:
                return f
        return formats[-1]

    def _generate_internal_regulation(self, index: int, extension: str, pages: int) -> dict:
        """Generate an internal regulation document with test:policy type."""
        templates = self.config["content_templates"]["internal_regulations"]
        titles = templates["titles"]
        departments = templates.get("departments", [{"id": "general", "name": "総務部"}])

        title = titles[index % len(titles)]
        department = random.choice(departments)
        policy_number = f"POL-{index + 1:03d}"
        policy_version = f"第{random.randint(1, 5)}版"

        # Generate dates
        effective_date = datetime(2024, 4, 1) + timedelta(days=random.randint(0, 365))
        review_date = effective_date + timedelta(days=365)

        if index >= len(titles):
            title = f"{title}_{index // len(titles) + 1}"

        content = self._create_regulation_content(title, extension, pages, policy_number, department["name"], effective_date, policy_version)

        return {
            "name": f"{title}.{extension}",
            "content": content,
            "properties": {
                "test:policyNumber": policy_number,
                "test:department": department["id"],
                "test:effectiveDate": effective_date.strftime('%Y-%m-%dT%H:%M:%S.000Z'),
                "test:reviewDate": review_date.strftime('%Y-%m-%dT%H:%M:%S.000Z'),
                "test:policyVersion": policy_version,
            },
        }

    def _create_regulation_content(self, title: str, extension: str, pages: int,
                                     policy_number: str = None, department: str = None,
                                     effective_date: datetime = None, policy_version: str = None) -> bytes:
        """Create regulation document content.

        Args:
            title: Document title
            extension: File extension (must be 'docx' or 'pptx')
            pages: Number of pages to generate
            policy_number: Policy document number
            department: Managing department name
            effective_date: Effective date of the policy
            policy_version: Version string of the policy

        Returns:
            Document content as bytes

        Raises:
            ValueError: If extension is not 'docx' or 'pptx'
        """
        if extension == "docx":
            return self._create_word_document(title, pages, policy_number, department, effective_date, policy_version)
        elif extension == "pptx":
            return self._create_powerpoint_document(title, pages)
        else:
            raise ValueError(
                f"Unsupported extension '{extension}' for internal regulations. "
                f"Only 'docx' and 'pptx' are supported. "
                f"Please check your config.yaml formats section."
            )

    def _create_word_document(self, title: str, pages: int,
                               policy_number: str = None, department: str = None,
                               effective_date: datetime = None, policy_version: str = None) -> bytes:
        """Create a Word document with regulation content."""
        doc = Document()

        doc.add_heading(title, 0)

        # Add policy metadata if provided
        if policy_number or department or effective_date or policy_version:
            if policy_number:
                doc.add_paragraph(f"規程番号: {policy_number}")
            if effective_date:
                doc.add_paragraph(f"施行日: {effective_date.strftime('%Y年%m月%d日')}")
            if department:
                doc.add_paragraph(f"管理部門: {department}")
            if policy_version:
                doc.add_paragraph(f"版数: {policy_version}")
            doc.add_paragraph("")

        doc.add_heading("第1章 総則", level=1)
        doc.add_paragraph(
            f"第1条（目的）\n"
            f"本規程は、当社における{title}に関する基本的事項を定め、"
            f"業務の適正かつ円滑な運営を図ることを目的とする。"
        )
        doc.add_paragraph(
            "第2条（適用範囲）\n"
            "本規程は、当社の全ての役員および従業員に適用する。"
        )
        doc.add_paragraph(
            "第3条（定義）\n"
            "本規程において使用する用語の定義は、以下のとおりとする。\n"
            "（1）「会社」とは、株式会社テスト商事をいう。\n"
            "（2）「従業員」とは、当社と雇用契約を締結した者をいう。\n"
            "（3）「所属長」とは、各部門の長をいう。"
        )

        for page in range(1, pages):
            doc.add_page_break()
            chapter_num = page + 1
            doc.add_heading(f"第{chapter_num}章 詳細規定", level=1)

            for section in range(1, 4):
                article_num = (page - 1) * 3 + section + 3
                doc.add_paragraph(
                    f"第{article_num}条（規定事項{section}）\n"
                    f"本条では、{title}に関する詳細な規定事項{section}について定める。\n"
                    f"1. 従業員は、本規程に従い、適切に業務を遂行しなければならない。\n"
                    f"2. 所属長は、部下の業務遂行状況を適切に管理・監督する責任を負う。\n"
                    f"3. 本条に違反した場合は、就業規則に基づき懲戒処分の対象となることがある。"
                )

        doc.add_heading("附則", level=1)
        effective_date_str = effective_date.strftime('%Y年%m月%d日') if effective_date else "令和6年4月1日"
        doc.add_paragraph(
            f"1. 本規程は、{effective_date_str}から施行する。\n"
            "2. 本規程の改廃は、取締役会の決議による。"
        )

        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue()

    def _create_powerpoint_document(self, title: str, pages: int) -> bytes:
        """Create a PowerPoint document with regulation content."""
        prs = Presentation()

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(2), PptxInches(9), PptxInches(2)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PptxPt(44)
        p.font.bold = True
        p.alignment = 1

        subtitle_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(4), PptxInches(9), PptxInches(1)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = "株式会社テスト商事"
        p.font.size = PptxPt(24)
        p.alignment = 1

        for page in range(1, pages):
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)

            header_box = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(0.5), PptxInches(9), PptxInches(1)
            )
            tf = header_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"第{page}章 規定内容"
            p.font.size = PptxPt(32)
            p.font.bold = True

            content_box = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(1.5), PptxInches(9), PptxInches(5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            points = [
                f"{title}の目的と適用範囲",
                "関係者の責任と義務",
                "具体的な手続きと方法",
                "違反時の対応について",
                "改定履歴と施行日",
            ]

            for i, point in enumerate(points):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {point}"
                p.font.size = PptxPt(20)
                p.space_before = PptxPt(12)

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def _generate_contract(self, index: int, pages: int) -> dict:
        """Generate a contract document (always PDF format) with rich content for RAG testing."""
        templates = self.config["content_templates"]["contracts"]
        titles = templates["titles"]
        companies = templates["company_names"]
        contract_types = templates.get("types", [{"id": "other", "name": "その他"}])

        title = titles[index % len(titles)]
        company = companies[index % len(companies)]
        contract_no = f"CTR-{datetime.now().year}-{index + 1:04d}"
        contract_type = random.choice(contract_types)

        # Generate contract dates
        start_date = datetime.now() - timedelta(days=random.randint(0, 90))
        end_date = start_date + timedelta(days=random.randint(365, 730))

        # Generate contract amount
        amount_ranges = templates.get("amount_ranges", [{"min": 1000000, "max": 10000000}])
        amount_range = random.choice(amount_ranges)
        contract_amount = random.randint(amount_range["min"], amount_range["max"])

        # Select payment term, confidentiality period, and penalty clause
        payment_terms = templates.get("payment_terms", ["契約締結後30日以内に一括払い"])
        confidentiality_periods = templates.get("confidentiality_periods", ["契約終了後3年間"])
        penalty_clauses = templates.get("penalty_clauses", ["契約金額の10%相当額を違約金として支払う"])

        payment_term = random.choice(payment_terms)
        confidentiality_period = random.choice(confidentiality_periods)
        penalty_clause = random.choice(penalty_clauses)

        # Create contract with rich content
        content = self._create_contract_pdf(
            title, company, contract_no, pages, start_date, end_date,
            contract_amount, payment_term, confidentiality_period, penalty_clause
        )

        return {
            "name": f"{title}_{contract_no}.pdf",
            "content": content,
            "properties": {
                "test:contractNumber": contract_no,
                "test:partyName": company,
                "test:startDate": start_date.strftime('%Y-%m-%dT%H:%M:%S.000Z'),
                "test:endDate": end_date.strftime('%Y-%m-%dT%H:%M:%S.000Z'),
                "test:contractType": contract_type["id"],
            },
        }

    def _create_contract_pdf(self, title: str, company: str, contract_no: str, pages: int,
                              start_date: datetime = None, end_date: datetime = None,
                              contract_amount: int = None, payment_term: str = None,
                              confidentiality_period: str = None, penalty_clause: str = None) -> bytes:
        """Create a PDF contract document with rich content for RAG analysis.

        Args:
            title: Contract title
            company: Counterparty company name
            contract_no: Contract number
            pages: Number of pages
            start_date: Contract start date
            end_date: Contract end date
            contract_amount: Contract amount in yen
            payment_term: Payment terms description
            confidentiality_period: Confidentiality period description
            penalty_clause: Penalty clause description

        Returns:
            PDF content as bytes
        """
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=2 * cm, bottomMargin=2 * cm)

        styles = getSampleStyleSheet()

        title_style = ParagraphStyle(
            "CustomTitle",
            parent=styles["Heading1"],
            fontName=self.JAPANESE_FONT,
            fontSize=18,
            alignment=1,
            spaceAfter=30,
        )

        body_style = ParagraphStyle(
            "CustomBody",
            parent=styles["Normal"],
            fontName=self.JAPANESE_FONT,
            fontSize=10,
            leading=14,
            spaceAfter=12,
        )

        # Use provided values or generate defaults
        if start_date is None:
            start_date = datetime.now() - timedelta(days=random.randint(1, 365))
        if end_date is None:
            end_date = start_date + timedelta(days=random.randint(365, 730))
        if contract_amount is None:
            contract_amount = random.randint(1000000, 50000000)
        if payment_term is None:
            payment_term = "契約締結後30日以内に一括払い"
        if confidentiality_period is None:
            confidentiality_period = "契約終了後3年間"
        if penalty_clause is None:
            penalty_clause = "契約金額の10%相当額を違約金として支払う"

        story = []

        # Title
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 20))

        # Contract metadata
        story.append(
            Paragraph(
                f"契約番号: {contract_no}",
                body_style,
            )
        )
        story.append(
            Paragraph(
                f"契約金額: ¥{contract_amount:,}（税別）",
                body_style,
            )
        )
        story.append(Spacer(1, 20))

        # Preamble
        story.append(
            Paragraph(
                f"株式会社テスト商事（以下「甲」という）と{company}（以下「乙」という）は、"
                f"以下のとおり{title}を締結する。",
                body_style,
            )
        )
        story.append(Spacer(1, 20))

        # Main articles with detailed content for RAG analysis
        articles = [
            ("第1条（目的）",
             f"本契約は、甲乙間における{title.replace('契約書', '')}に関する基本的事項を定めることを目的とする。"
             f"本契約に基づく業務の範囲、対価、履行方法等については、本契約および別途締結する個別契約に定めるところによる。"),

            ("第2条（契約期間）",
             f"本契約の有効期間は、{start_date.strftime('%Y年%m月%d日')}から{end_date.strftime('%Y年%m月%d日')}までとする。"
             "ただし、期間満了の1ヶ月前までに甲乙いずれからも書面による解約の申し出がない場合は、"
             "同一条件で更に1年間延長されるものとし、以後も同様とする。"),

            ("第3条（契約金額および支払条件）",
             f"本契約に基づく契約金額は、金{contract_amount:,}円（消費税別）とする。"
             f"支払条件は、{payment_term}とし、甲が乙に指定する銀行口座への振込みにより支払うものとする。"
             "振込手数料は甲の負担とする。"),

            ("第4条（秘密保持）",
             "甲および乙は、本契約に関連して知り得た相手方の技術上または営業上の秘密情報"
             "（以下「秘密情報」という）を、相手方の書面による事前の承諾なく第三者に開示または漏洩してはならない。"
             f"この秘密保持義務は、本契約終了後も{confidentiality_period}存続するものとする。"
             "なお、以下の情報は秘密情報に含まれないものとする。"
             "（1）開示を受けた時点で既に公知であった情報、"
             "（2）開示を受けた後、受領者の責によらずに公知となった情報、"
             "（3）開示を受けた時点で既に受領者が保有していた情報、"
             "（4）正当な権限を有する第三者から秘密保持義務を負うことなく適法に取得した情報。"),

            ("第5条（知的財産権）",
             "本契約に基づき甲が乙に提供する成果物に関する著作権（著作権法第27条および第28条の権利を含む）"
             "その他一切の知的財産権は、対価の完済をもって甲に帰属するものとする。"
             "乙は、甲に対し、当該成果物に関して著作者人格権を行使しないものとする。"),

            ("第6条（損害賠償）",
             "甲または乙が本契約に違反し、相手方に損害を与えた場合は、{penalty_clause}ものとする。"
             .replace("{penalty_clause}", penalty_clause) +
             "ただし、当該違反が天災地変その他不可抗力による場合はこの限りでない。"),

            ("第7条（契約解除）",
             "甲または乙は、相手方が以下の各号のいずれかに該当した場合、何らの催告なく直ちに本契約を解除することができる。"
             "（1）本契約に違反し、相当の期間を定めて催告したにもかかわらず是正されない場合、"
             "（2）支払停止または支払不能となった場合、"
             "（3）手形または小切手が不渡りとなった場合、"
             "（4）差押え、仮差押え、仮処分または競売の申立てがあった場合、"
             "（5）破産手続開始、民事再生手続開始、会社更生手続開始または特別清算開始の申立てがあった場合、"
             "（6）解散または事業の全部もしくは重要な一部を第三者に譲渡した場合。"),

            ("第8条（反社会的勢力の排除）",
             "甲および乙は、自己または自己の役員が暴力団、暴力団員、暴力団準構成員、暴力団関係企業、"
             "総会屋、社会運動等標榜ゴロ、特殊知能暴力集団その他これらに準ずる反社会的勢力"
             "（以下「反社会的勢力」という）に該当しないことを表明し、将来にわたっても該当しないことを確約する。"
             "甲または乙が本条に違反した場合、相手方は何らの催告なく直ちに本契約を解除できるものとする。"),

            ("第9条（権利義務の譲渡禁止）",
             "甲および乙は、相手方の書面による事前の承諾なく、本契約上の地位または本契約に基づく権利義務の"
             "全部もしくは一部を第三者に譲渡し、または担保に供してはならない。"),

            ("第10条（協議事項）",
             "本契約に定めのない事項または本契約の条項の解釈に疑義が生じた場合は、"
             "甲乙誠意をもって協議し、解決するものとする。"),

            ("第11条（管轄裁判所）",
             "本契約に関する一切の紛争については、東京地方裁判所を第一審の専属的合意管轄裁判所とする。"),
        ]

        for article_title, article_content in articles:
            story.append(Paragraph(f"<b>{article_title}</b>", body_style))
            story.append(Paragraph(article_content, body_style))
            story.append(Spacer(1, 10))

        # Additional pages with supplementary articles
        for page in range(1, pages - 1):
            story.append(Spacer(1, 50))

            supplementary_articles = [
                (f"第{12 + (page - 1) * 2}条（報告義務）",
                 f"乙は、本契約に基づく業務の進捗状況について、甲の求めに応じて速やかに報告するものとする。"
                 f"また、乙は、業務遂行上重大な問題が発生した場合は、直ちに甲に報告し、その指示を仰ぐものとする。"),

                (f"第{13 + (page - 1) * 2}条（再委託）",
                 f"乙は、甲の書面による事前の承諾を得た場合に限り、本契約に基づく業務の全部または一部を第三者に再委託することができる。"
                 f"この場合、乙は当該第三者に対し、本契約に基づき乙が負う義務と同等の義務を負わせるものとし、"
                 f"当該第三者の行為について甲に対し一切の責任を負うものとする。"),
            ]

            for article_title, article_content in supplementary_articles:
                story.append(Paragraph(f"<b>{article_title}</b>", body_style))
                story.append(Paragraph(article_content, body_style))
                story.append(Spacer(1, 10))

        # Closing statement
        story.append(Spacer(1, 30))
        story.append(
            Paragraph(
                "本契約締結の証として、本書2通を作成し、甲乙記名押印の上、各1通を保有する。",
                body_style,
            )
        )
        story.append(Spacer(1, 20))

        # Contract summary
        story.append(
            Paragraph(
                f"契約締結日: {start_date.strftime('%Y年%m月%d日')}",
                body_style,
            )
        )
        story.append(
            Paragraph(
                f"契約期間: {start_date.strftime('%Y年%m月%d日')} ～ {end_date.strftime('%Y年%m月%d日')}",
                body_style,
            )
        )
        story.append(
            Paragraph(
                f"契約金額: ¥{contract_amount:,}（税別）",
                body_style,
            )
        )
        story.append(Spacer(1, 30))

        # Signature blocks
        story.append(Paragraph("<b>甲（発注者）:</b>", body_style))
        story.append(Paragraph("  株式会社テスト商事", body_style))
        story.append(Paragraph("  〒100-0001 東京都千代田区丸の内1-1-1", body_style))
        story.append(Paragraph("  代表取締役 山田 太郎", body_style))
        story.append(Spacer(1, 20))

        story.append(Paragraph("<b>乙（受注者）:</b>", body_style))
        story.append(Paragraph(f"  {company}", body_style))
        story.append(Paragraph("  〒150-0001 東京都渋谷区神宮前1-2-3", body_style))
        story.append(Paragraph("  代表取締役 鈴木 一郎", body_style))

        doc.build(story)
        return buffer.getvalue()

    def _generate_invoice(self, index: int, extension: str) -> dict:
        """Generate an invoice document."""
        templates = self.config["content_templates"]["invoices"]
        items = templates["items"]
        customer_names = templates.get("customer_names", ["株式会社サンプル顧客"])
        statuses = templates.get("statuses", [{"id": "draft", "name": "未送付"}])

        invoice_no = f"INV-{datetime.now().year}-{index + 1:05d}"
        customer_name = random.choice(customer_names)
        status = random.choice(statuses)
        issue_date = datetime.now() - timedelta(days=random.randint(0, 180))

        num_items = random.randint(1, 5)
        selected_items = random.sample(items, min(num_items, len(items)))

        total_sum = 0
        invoice_items = []
        for item in selected_items:
            quantity = random.randint(1, 10)
            unit_price = item["unit_price"]
            amount = quantity * unit_price
            total_sum += amount
            invoice_items.append(
                {
                    "name": item["name"],
                    "quantity": quantity,
                    "unit_price": unit_price,
                    "amount": amount,
                }
            )

        content = self._create_invoice_excel(invoice_no, invoice_items, total_sum, templates, customer_name, issue_date)

        return {
            "name": f"請求書_{invoice_no}.xlsx",
            "content": content,
            "properties": {
                "test:invoiceNumber": invoice_no,
                "test:customerName": customer_name,
                "test:amount": total_sum,
                "test:issueDate": issue_date.strftime('%Y-%m-%dT%H:%M:%S.000Z'),
                "test:invoiceStatus": status["id"],
            },
        }

    def _create_invoice_excel(
        self, invoice_no: str, items: list, total_sum: int, templates: dict,
        customer_name: str = None, issue_date: datetime = None
    ) -> bytes:
        """Create an Excel invoice document."""
        wb = Workbook()
        ws = wb.active
        ws.title = "請求書"

        ws.column_dimensions["A"].width = 5
        ws.column_dimensions["B"].width = 40
        ws.column_dimensions["C"].width = 10
        ws.column_dimensions["D"].width = 15
        ws.column_dimensions["E"].width = 18

        title_font = Font(size=20, bold=True)
        header_font = Font(size=11, bold=True)
        header_fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")
        border = Border(
            left=Side(style="thin"),
            right=Side(style="thin"),
            top=Side(style="thin"),
            bottom=Side(style="thin"),
        )

        ws.merge_cells("A1:E1")
        ws["A1"] = "請 求 書"
        ws["A1"].font = title_font
        ws["A1"].alignment = Alignment(horizontal="center")

        ws["A3"] = f"請求書番号: {invoice_no}"
        invoice_date = issue_date if issue_date else datetime.now() - timedelta(days=random.randint(0, 30))
        ws["D3"] = f"発行日: {invoice_date.strftime('%Y年%m月%d日')}"

        ws["A5"] = "御中"
        ws["A6"] = customer_name if customer_name else f"株式会社サンプル顧客{random.randint(1, 100)}"

        ws["D5"] = templates["company_name"]
        ws["D6"] = templates["company_address"]
        ws["D7"] = f"TEL: {templates['company_tel']}"

        ws["A9"] = "下記の通りご請求申し上げます。"

        ws.merge_cells("A11:E11")
        ws["A11"] = f"ご請求金額: ¥{total_sum:,}（税込）"
        ws["A11"].font = Font(size=14, bold=True)
        ws["A11"].alignment = Alignment(horizontal="center")

        headers = ["No.", "品目", "数量", "単価", "金額"]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=13, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.border = border
            cell.alignment = Alignment(horizontal="center")

        for row_idx, item in enumerate(items, 14):
            ws.cell(row=row_idx, column=1, value=row_idx - 13).border = border
            ws.cell(row=row_idx, column=2, value=item["name"]).border = border
            ws.cell(row=row_idx, column=3, value=item["quantity"]).border = border
            ws.cell(row=row_idx, column=4, value=f"¥{item['unit_price']:,}").border = border
            ws.cell(row=row_idx, column=5, value=f"¥{item['amount']:,}").border = border

            for col in range(1, 6):
                ws.cell(row=row_idx, column=col).alignment = Alignment(horizontal="center" if col != 2 else "left")

        subtotal_row = 14 + len(items)
        ws.merge_cells(f"A{subtotal_row}:D{subtotal_row}")
        ws.cell(row=subtotal_row, column=1, value="小計").border = border
        ws.cell(row=subtotal_row, column=5, value=f"¥{total_sum:,}").border = border

        tax = int(total_sum * 0.1)
        tax_row = subtotal_row + 1
        ws.merge_cells(f"A{tax_row}:D{tax_row}")
        ws.cell(row=tax_row, column=1, value="消費税（10%）").border = border
        ws.cell(row=tax_row, column=5, value=f"¥{tax:,}").border = border

        total_row = tax_row + 1
        ws.merge_cells(f"A{total_row}:D{total_row}")
        ws.cell(row=total_row, column=1, value="合計").border = border
        ws.cell(row=total_row, column=1).font = header_font
        ws.cell(row=total_row, column=5, value=f"¥{total_sum + tax:,}").border = border
        ws.cell(row=total_row, column=5).font = header_font

        note_row = total_row + 2
        ws[f"A{note_row}"] = "備考:"
        ws[f"A{note_row + 1}"] = "・お支払期限: 請求書発行日より30日以内"
        ws[f"A{note_row + 2}"] = "・振込先: テスト銀行 本店 普通 1234567"

        buffer = io.BytesIO()
        wb.save(buffer)
        return buffer.getvalue()

    def _upload_document(
        self,
        folder_id: str,
        name: str,
        content: bytes,
        content_type: str,
        properties: dict,
    ) -> None:
        """Upload a document to the repository using Browser Binding."""
        try:
            browser_url = self._get_browser_url()
            mime_type = self._get_mime_type(name)

            # Prepare multipart form data
            files = {
                "content": (name, io.BytesIO(content), mime_type)
            }

            data = {
                "cmisaction": "createDocument",
                "folderId": folder_id,
                "propertyId[0]": "cmis:objectTypeId",
                "propertyValue[0]": content_type,
                "propertyId[1]": "cmis:name",
                "propertyValue[1]": name,
            }

            # Add custom properties
            prop_idx = 2
            for key, value in properties.items():
                data[f"propertyId[{prop_idx}]"] = key
                # Convert values to string format expected by CMIS Browser Binding
                if isinstance(value, bool):
                    data[f"propertyValue[{prop_idx}]"] = str(value).lower()
                elif isinstance(value, int):
                    # Integer properties must be passed as string
                    data[f"propertyValue[{prop_idx}]"] = str(value)
                else:
                    data[f"propertyValue[{prop_idx}]"] = value
                prop_idx += 1

            response = requests.post(
                browser_url,
                auth=self._get_cmis_auth(),
                data=data,
                files=files,
                timeout=HTTP_TIMEOUT
            )

            # Check for successful upload
            # Browser Binding returns 200/201 with JSON containing objectId on success
            if response.status_code in (200, 201):
                try:
                    result = response.json()
                    # Success: response contains properties or succinctProperties with objectId
                    if "properties" in result or "succinctProperties" in result:
                        return  # Upload successful
                except Exception:
                    pass

            # Check for error responses
            if "exception" in response.text:
                # Skip "already exists" errors silently
                if "already exists" in response.text.lower():
                    return
                print(f"      Error uploading {name}: {response.text[:200]}")

        except Exception as e:
            if "already exists" in str(e).lower():
                pass
            else:
                print(f"      Error uploading {name}: {e}")

    def _get_mime_type(self, filename: str) -> str:
        """Get MIME type from filename.

        Uses a predefined map for Office formats (which mimetypes module may not know),
        then falls back to the mimetypes module, and finally to application/octet-stream.
        """
        ext = filename.lower().split(".")[-1]
        # Office formats that mimetypes module may not recognize correctly
        office_mime_types = {
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "pdf": "application/pdf",
        }
        if ext in office_mime_types:
            return office_mime_types[ext]

        # Fall back to mimetypes module for other extensions
        mime_type, _ = mimetypes.guess_type(filename)
        return mime_type if mime_type else "application/octet-stream"

    def run(self) -> None:
        """Run the complete test environment setup."""
        print("=" * 60)
        print("NemakiWare Test Environment Setup")
        print("=" * 60)

        try:
            self.connect_cmis()

            self.create_groups()

            self.create_users()

            self.register_custom_types()

            self.create_organization_folders()

            self.create_root_folders()

            print("\n" + "=" * 60)
            print("Test environment setup completed successfully!")
            print("=" * 60)

        except Exception as e:
            print(f"\nError during setup: {e}")
            raise


def main():
    """Main entry point."""
    parser = argparse.ArgumentParser(
        description="Set up NemakiWare test environment with sample data"
    )
    parser.add_argument(
        "--config",
        "-c",
        default="config.yaml",
        help="Path to configuration file (default: config.yaml)",
    )
    parser.add_argument(
        "--seed",
        "-s",
        type=int,
        default=None,
        help="Random seed for reproducible test data generation (default: None)",
    )
    args = parser.parse_args()

    # Apply random seed from argument or environment variable
    seed = args.seed
    if seed is None and os.environ.get("NEMAKI_RANDOM_SEED"):
        try:
            seed = int(os.environ["NEMAKI_RANDOM_SEED"])
        except ValueError:
            print("Warning: NEMAKI_RANDOM_SEED is not a valid integer, ignoring")

    if seed is not None:
        random.seed(seed)
        print(f"Using random seed: {seed}")

    if not os.path.exists(args.config):
        print(f"Error: Configuration file not found: {args.config}")
        sys.exit(1)

    setup = TestEnvironmentSetup(args.config)
    setup.run()


if __name__ == "__main__":
    main()
